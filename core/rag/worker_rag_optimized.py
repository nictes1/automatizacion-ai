#!/usr/bin/env python3
"""
PulpoAI RAG Worker Optimizado
Sistema de búsqueda semántica usando Ollama + GPU
"""

import os
import json
import hashlib
import asyncio
import logging
from typing import List, Dict, Any, Optional
from pathlib import Path
import uuid
from datetime import datetime
from contextlib import asynccontextmanager

# FastAPI
from fastapi import FastAPI, File, UploadFile, HTTPException, Header, Depends
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import uvicorn

# Base de datos
import psycopg2
from psycopg2.extras import RealDictCursor
import redis

# Weaviate
import weaviate

# HTTP requests para Ollama
import requests

# Sistema de caché
from embedding_cache import embedding_cache

# Sistema de monitoreo
from monitoring import system_monitor, get_system_health

# Procesador inteligente de documentos
from smart_document_processor import SmartDocumentProcessor

# Procesamiento de documentos
import pypdf
from docx import Document
import openpyxl
from pptx import Presentation
import pytesseract
import magic

# Utilidades
import tiktoken
from dotenv import load_dotenv

# Configurar logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Cargar variables de entorno
load_dotenv()

# Configuración
DATABASE_URL = os.getenv("DATABASE_URL", "postgresql://pulpo:pulpo@localhost:5432/pulpo")
WEAVIATE_URL = os.getenv("WEAVIATE_URL", "http://localhost:8080")
REDIS_URL = os.getenv("REDIS_URL", "redis://localhost:6379")
WORKER_PORT = int(os.getenv("WORKER_PORT", 8002))
OLLAMA_URL = os.getenv("OLLAMA_URL", "http://ollama:11434")

# Inicializar FastAPI
app = FastAPI(
    title="PulpoAI RAG Worker Optimizado",
    description="Sistema de búsqueda semántica usando Ollama + GPU",
    version="2.0.0"
)

# CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Modelos Pydantic
class SearchRequest(BaseModel):
    query: str
    limit: int = 10
    similarity_threshold: float = 0.7
    include_metadata: bool = True

class IngestRequest(BaseModel):
    title: Optional[str] = None
    metadata: Optional[Dict[str, Any]] = {}

class SearchResponse(BaseModel):
    results: List[Dict[str, Any]]
    total: int
    query: str
    processing_time_ms: float

class IngestResponse(BaseModel):
    document_id: str
    chunks_created: int
    status: str
    processing_time_ms: float

class RAGSystem:
    def __init__(self):
        self.pg_conn = None
        self.weaviate_client = None
        self.redis_client = None
        self.chunk_tokenizer = None
        self.ollama_url = OLLAMA_URL
        self.smart_processor = SmartDocumentProcessor(OLLAMA_URL, "llama3.1:8b")
        
    async def initialize(self):
        """Inicializar conexiones y modelos"""
        try:
            # PostgreSQL
            self.pg_conn = psycopg2.connect(DATABASE_URL)
            logger.info("✅ PostgreSQL conectado")
            
            # Weaviate
            self.weaviate_client = weaviate.Client(WEAVIATE_URL)
            await self._setup_weaviate_schema()
            logger.info("✅ Weaviate conectado")
            
            # Redis
            self.redis_client = redis.from_url(REDIS_URL)
            self.redis_client.ping()
            logger.info("✅ Redis conectado")
            
            # Tokenizer para chunking
            self.chunk_tokenizer = tiktoken.get_encoding("cl100k_base")
            logger.info("✅ Tokenizer cargado")
            
            # Verificar Ollama
            await self._verify_ollama()
            logger.info("✅ Ollama verificado")
            
        except Exception as e:
            logger.error(f"❌ Error inicializando RAG System: {e}")
            raise

    async def _verify_ollama(self):
        """Verificar que Ollama esté funcionando"""
        try:
            response = requests.get(f"{self.ollama_url}/api/tags", timeout=10)
            if response.status_code == 200:
                models = response.json().get("models", [])
                logger.info(f"✅ Ollama disponible con {len(models)} modelos")
            else:
                raise Exception(f"Ollama no responde: {response.status_code}")
        except Exception as e:
            logger.error(f"❌ Error verificando Ollama: {e}")
            raise

    async def _setup_weaviate_schema(self):
        """Configurar esquema de Weaviate"""
        try:
            # Verificar si ya existe la clase
            if self.weaviate_client.schema.exists("Document"):
                logger.info("✅ Esquema Weaviate ya existe")
                return
            
            # Crear esquema
            schema = {
                "class": "Document",
                "description": "Documentos del sistema RAG",
                "vectorizer": "none",  # Usamos embeddings externos
                "properties": [
                    {
                        "name": "workspace_id",
                        "dataType": ["string"],
                        "description": "ID del workspace"
                    },
                    {
                        "name": "document_id",
                        "dataType": ["string"],
                        "description": "ID del documento"
                    },
                    {
                        "name": "chunk_index",
                        "dataType": ["int"],
                        "description": "Índice del chunk"
                    },
                    {
                        "name": "content",
                        "dataType": ["text"],
                        "description": "Contenido del chunk"
                    },
                    {
                        "name": "title",
                        "dataType": ["string"],
                        "description": "Título del documento"
                    },
                    {
                        "name": "file_path",
                        "dataType": ["string"],
                        "description": "Ruta del archivo"
                    },
                    {
                        "name": "metadata",
                        "dataType": ["text"],
                        "description": "Metadatos del documento"
                    }
                ]
            }
            
            self.weaviate_client.schema.create_class(schema)
            logger.info("✅ Esquema Weaviate creado")
            
        except Exception as e:
            logger.warning(f"⚠️ Esquema Weaviate puede que ya exista: {e}")
            # No fallar si ya existe

    def extract_text_from_file(self, file_path: str) -> str:
        """Extraer texto de diferentes tipos de archivos"""
        try:
            file_ext = Path(file_path).suffix.lower()
            mime_type = magic.from_file(file_path, mime=True)
            
            if file_ext == '.pdf':
                return self._extract_pdf_text(file_path)
            elif file_ext in ['.docx', '.doc']:
                return self._extract_docx_text(file_path)
            elif file_ext in ['.xlsx', '.xls']:
                return self._extract_excel_text(file_path)
            elif file_ext in ['.pptx', '.ppt']:
                return self._extract_pptx_text(file_path)
            elif file_ext in ['.txt', '.md']:
                return self._extract_text_file(file_path)
            elif mime_type.startswith('image/'):
                return self._extract_image_text(file_path)
            else:
                # Intentar como texto plano
                return self._extract_text_file(file_path)
                
        except Exception as e:
            logger.error(f"❌ Error extrayendo texto de {file_path}: {e}")
            return ""

    def _extract_pdf_text(self, file_path: str) -> str:
        """Extraer texto de PDF"""
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = pypdf.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text

    def _extract_docx_text(self, file_path: str) -> str:
        """Extraer texto de DOCX"""
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text

    def _extract_excel_text(self, file_path: str) -> str:
        """Extraer texto de Excel"""
        workbook = openpyxl.load_workbook(file_path)
        text = ""
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            for row in sheet.iter_rows():
                row_text = " ".join([str(cell.value) for cell in row if cell.value])
                if row_text:
                    text += row_text + "\n"
        return text

    def _extract_pptx_text(self, file_path: str) -> str:
        """Extraer texto de PowerPoint"""
        presentation = Presentation(file_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def _extract_text_file(self, file_path: str) -> str:
        """Extraer texto de archivo de texto"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            with open(file_path, 'r', encoding='latin-1') as file:
                return file.read()

    def _extract_image_text(self, file_path: str) -> str:
        """Extraer texto de imagen usando OCR"""
        try:
            return pytesseract.image_to_string(file_path)
        except Exception as e:
            logger.error(f"❌ Error en OCR: {e}")
            return ""

    def chunk_text(self, text: str, chunk_size: int = 200, overlap: int = 50) -> List[str]:
        """Dividir texto en chunks"""
        if not text.strip():
            return []
        
        # Tokenizar
        tokens = self.chunk_tokenizer.encode(text)
        
        chunks = []
        for i in range(0, len(tokens), chunk_size - overlap):
            chunk_tokens = tokens[i:i + chunk_size]
            chunk_text = self.chunk_tokenizer.decode(chunk_tokens)
            chunks.append(chunk_text.strip())
        
        return [chunk for chunk in chunks if chunk]

    async def generate_embeddings_ollama(self, texts: List[str]) -> List[List[float]]:
        """Generar embeddings usando Ollama con caché"""
        try:
            model_name = "nomic-embed-text:latest"
            embeddings = []
            texts_to_generate = []
            text_indices = []
            
            # Verificar caché para todos los textos
            for i, text in enumerate(texts):
                cached_embedding = embedding_cache.get_embedding(text, model_name)
                if cached_embedding:
                    embeddings.append(cached_embedding)
                    logger.debug(f"✅ Embedding desde caché para: {text[:50]}...")
                else:
                    embeddings.append(None)  # Placeholder
                    texts_to_generate.append(text)
                    text_indices.append(i)
            
            # Generar embeddings faltantes
            if texts_to_generate:
                logger.info(f"🔄 Generando {len(texts_to_generate)} embeddings nuevos...")
                
                for text in texts_to_generate:
                    response = requests.post(
                        f"{self.ollama_url}/api/embeddings",
                        json={
                            "model": model_name,
                            "prompt": text
                        },
                        timeout=30
                    )
                    
                    if response.status_code == 200:
                        embedding = response.json()["embedding"]
                        # Almacenar en caché
                        embedding_cache.set_embedding(text, model_name, embedding)
                        # Encontrar índice y asignar
                        idx = texts_to_generate.index(text)
                        original_idx = text_indices[idx]
                        embeddings[original_idx] = embedding
                        logger.debug(f"✅ Embedding generado y cacheado para: {text[:50]}...")
                    else:
                        logger.error(f"❌ Error generando embedding: {response.status_code}")
                        return []
            
            return embeddings
            
        except Exception as e:
            logger.error(f"❌ Error generando embeddings con Ollama: {e}")
            return []

    async def save_document_to_pg(self, workspace_id: str, title: str, file_path: str, 
                                 metadata: Dict[str, Any], text: str) -> str:
        """Guardar documento en PostgreSQL"""
        try:
            # Calcular hash del archivo
            with open(file_path, 'rb') as f:
                file_hash = hashlib.sha256(f.read()).hexdigest()
            
            # Primero insertar en la tabla files
            cursor = self.pg_conn.cursor()
            cursor.execute("""
                INSERT INTO pulpo.files (workspace_id, storage_uri, filename, mime_type, sha256, bytes, status)
                VALUES (%s, %s, %s, %s, %s, %s, %s)
                RETURNING id
            """, (workspace_id, str(file_path), title, "text/plain", file_hash, os.path.getsize(file_path), "processed"))
            
            file_id = cursor.fetchone()[0]
            
            # Luego insertar en la tabla documents
            token_count = len(self.chunk_tokenizer.encode(text))
            cursor.execute("""
                INSERT INTO pulpo.documents (workspace_id, file_id, title, language, raw_text, token_count)
                VALUES (%s, %s, %s, %s, %s, %s)
                RETURNING id
            """, (workspace_id, file_id, title, "es", text, token_count))
            
            document_id = cursor.fetchone()[0]
            self.pg_conn.commit()
            cursor.close()
            
            return str(document_id)
            
        except Exception as e:
            logger.error(f"❌ Error guardando documento en PostgreSQL: {e}")
            raise

    async def save_chunks_to_weaviate(self, workspace_id: str, document_id: str, 
                                    chunks: List[str], embeddings: List[List[float]], 
                                    title: str, file_path: str, metadata: Dict[str, Any]):
        """Guardar chunks en Weaviate"""
        try:
            for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
                chunk_data = {
                    "workspace_id": workspace_id,
                    "document_id": document_id,
                    "chunk_index": i,
                    "content": chunk,
                    "title": title,
                    "file_path": file_path,
                    "metadata": json.dumps(metadata)
                }
                
                self.weaviate_client.data_object.create(
                    data_object=chunk_data,
                    class_name="Document",
                    vector=embedding
                )
            
            logger.info(f"✅ {len(chunks)} chunks guardados en Weaviate")
            
        except Exception as e:
            logger.error(f"❌ Error guardando chunks en Weaviate: {e}")
            raise

    async def search_similar(self, query: str, workspace_id: str, limit: int = 10, 
                           similarity_threshold: float = 0.7) -> List[Dict[str, Any]]:
        """Buscar documentos similares"""
        try:
            # Generar embedding de la consulta
            query_embeddings = await self.generate_embeddings_ollama([query])
            if not query_embeddings:
                return []
            
            query_vector = query_embeddings[0]
            
            # Buscar en Weaviate
            result = self.weaviate_client.query.get(
                "Document", 
                ["workspace_id", "document_id", "chunk_index", "content", "title", "file_path", "metadata"]
            ).with_near_vector({
                "vector": query_vector
            }).with_where({
                "path": ["workspace_id"],
                "operator": "Equal",
                "valueString": workspace_id
            }).with_limit(limit).do()
            
            results = []
            if "data" in result and "Get" in result["data"]:
                for item in result["data"]["Get"]["Document"]:
                    results.append({
                        "content": item["content"],
                        "title": item["title"],
                        "file_path": item["file_path"],
                        "chunk_index": item["chunk_index"],
                        "metadata": item["metadata"],
                        "similarity": 1.0  # Weaviate no devuelve score directamente
                    })
            
            return results
            
        except Exception as e:
            logger.error(f"❌ Error en búsqueda: {e}")
            return []

# Instancia global del sistema RAG
rag_system = RAGSystem()

@asynccontextmanager
async def lifespan(app: FastAPI):
    """Gestión del ciclo de vida de la aplicación"""
    await rag_system.initialize()
    yield
    if rag_system.pg_conn:
        rag_system.pg_conn.close()

app.router.lifespan_context = lifespan

@app.get("/health")
async def health_check():
    """Health check"""
    return {"status": "healthy", "service": "pulpo-rag-optimized"}

@app.get("/cache/stats")
async def get_cache_stats():
    """Obtener estadísticas del caché de embeddings"""
    try:
        stats = embedding_cache.get_cache_stats()
        return {
            "status": "success",
            "cache_stats": stats,
            "cache_healthy": embedding_cache.is_healthy()
        }
    except Exception as e:
        logger.error(f"❌ Error obteniendo estadísticas del caché: {e}")
        return {"status": "error", "message": str(e)}

@app.post("/cache/clear")
async def clear_cache():
    """Limpiar caché de embeddings"""
    try:
        deleted_count = embedding_cache.clear_cache()
        return {
            "status": "success",
            "deleted_embeddings": deleted_count,
            "message": f"Se eliminaron {deleted_count} embeddings del caché"
        }
    except Exception as e:
        logger.error(f"❌ Error limpiando caché: {e}")
        return {"status": "error", "message": str(e)}

@app.get("/monitoring/health")
async def get_monitoring_health():
    """Obtener resumen de salud del sistema"""
    try:
        health_summary = get_system_health()
        return {
            "status": "success",
            "health": health_summary
        }
    except Exception as e:
        logger.error(f"❌ Error obteniendo salud del sistema: {e}")
        return {"status": "error", "message": str(e)}

@app.get("/monitoring/metrics")
async def get_monitoring_metrics():
    """Obtener métricas detalladas del sistema"""
    try:
        system_metrics = system_monitor.get_system_metrics()
        service_metrics = await system_monitor.check_all_services()
        
        return {
            "status": "success",
            "system_metrics": {
                "timestamp": system_metrics.timestamp.isoformat(),
                "cpu_percent": system_metrics.cpu_percent,
                "memory_percent": system_metrics.memory_percent,
                "memory_used_mb": system_metrics.memory_used_mb,
                "disk_usage_percent": system_metrics.disk_usage_percent,
                "gpu_memory_used_mb": system_metrics.gpu_memory_used_mb,
                "gpu_utilization_percent": system_metrics.gpu_utilization_percent
            },
            "service_metrics": {
                service_name: {
                    "status": metrics.status,
                    "response_time_ms": metrics.response_time_ms,
                    "error_count": metrics.error_count,
                    "timestamp": metrics.timestamp.isoformat()
                }
                for service_name, metrics in service_metrics.items()
            }
        }
    except Exception as e:
        logger.error(f"❌ Error obteniendo métricas: {e}")
        return {"status": "error", "message": str(e)}

@app.post("/rag/ingest", response_model=IngestResponse)
async def ingest_document(
    file: UploadFile = File(...),
    workspace_id: str = Header(..., alias="x-workspace-id"),
    title: Optional[str] = None,
    metadata: Optional[str] = None
):
    """Ingerir un documento al sistema RAG"""
    start_time = datetime.now()
    
    try:
        # Validar archivo
        if not file.filename:
            raise HTTPException(status_code=400, detail="No se proporcionó archivo")
        
        # Crear directorio de uploads si no existe
        upload_dir = Path("uploads")
        upload_dir.mkdir(exist_ok=True)
        
        # Guardar archivo temporalmente
        file_path = upload_dir / f"{uuid.uuid4()}_{file.filename}"
        with open(file_path, "wb") as buffer:
            content = await file.read()
            buffer.write(content)
        
        # Procesar metadatos
        doc_metadata = {}
        if metadata:
            try:
                doc_metadata = json.loads(metadata)
            except:
                pass
        
        # Extraer texto
        text = rag_system.extract_text_from_file(str(file_path))
        if not text.strip():
            raise HTTPException(status_code=400, detail="No se pudo extraer texto del archivo")
        
        # Procesamiento inteligente del documento
        doc_title = title or file.filename
        if "menu" in doc_title.lower() or "restaurante" in doc_title.lower():
            # Usar procesador inteligente para menús
            smart_chunks = rag_system.smart_processor.process_document(text, "menu")
            chunks = [chunk["content"] for chunk in smart_chunks]
            chunk_metadata = [chunk["metadata"] for chunk in smart_chunks]
        else:
            # Chunking tradicional para otros documentos
            chunks = rag_system.chunk_text(text)
            chunk_metadata = [{"type": "general"}] * len(chunks)
        
        if not chunks:
            raise HTTPException(status_code=400, detail="No se pudieron crear chunks del texto")
        
        # Generar embeddings usando Ollama
        embeddings = await rag_system.generate_embeddings_ollama(chunks)
        if not embeddings:
            raise HTTPException(status_code=500, detail="Error generando embeddings")
        
        # Guardar documento en PostgreSQL
        doc_title = title or file.filename
        document_id = await rag_system.save_document_to_pg(
            workspace_id, doc_title, str(file_path), doc_metadata, text
        )
        
        # Guardar chunks en Weaviate con metadatos enriquecidos
        for i, (chunk, embedding, metadata) in enumerate(zip(chunks, embeddings, chunk_metadata)):
            await rag_system.save_chunks_to_weaviate(
                workspace_id, document_id, [chunk], [embedding], doc_title, str(file_path), metadata
            )
        
        # Limpiar archivo temporal
        file_path.unlink()
        
        processing_time = (datetime.now() - start_time).total_seconds() * 1000
        
        return IngestResponse(
            document_id=document_id,
            chunks_created=len(chunks),
            status="success",
            processing_time_ms=processing_time
        )
        
    except Exception as e:
        logger.error(f"❌ Error en ingesta: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/rag/search", response_model=SearchResponse)
async def search_documents(
    request: SearchRequest,
    workspace_id: str = Header(..., alias="x-workspace-id")
):
    """Buscar documentos similares"""
    start_time = datetime.now()
    
    try:
        results = await rag_system.search_similar(
            request.query, workspace_id, request.limit, request.similarity_threshold
        )
        
        processing_time = (datetime.now() - start_time).total_seconds() * 1000
        
        return SearchResponse(
            results=results,
            total=len(results),
            query=request.query,
            processing_time_ms=processing_time
        )
        
    except Exception as e:
        logger.error(f"❌ Error en búsqueda: {e}")
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=WORKER_PORT)
